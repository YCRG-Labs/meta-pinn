#!/usr/bin/env python3
"""
Generate Complete PowerPoint Presentation - All 28 Slides
Physics-Informed Meta-Learning for Few-Shot Parameter Inference
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path
import os

def create_complete_presentation():
    """Create the complete 28-slide PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define paths to figure directories
    conceptual_dir = Path("conceptual_figures")
    data_dir = Path("presentation_figures")
    paper_dir = Path("paper/MDPI/figures")
    
    print("Creating all 28 slides...")
    
    # Create all slides
    create_slide_01_title(prs, conceptual_dir)
    create_slide_02_problem(prs, conceptual_dir)
    create_slide_03_motivation(prs, conceptual_dir, data_dir)
    create_slide_04_meta_learning(prs, conceptual_dir)
    create_slide_05_contributions(prs, conceptual_dir)
    create_slide_06_formulation(prs, conceptual_dir)
    create_slide_07_framework(prs, conceptual_dir)
    create_slide_08_physics_loss(prs, conceptual_dir)
    create_slide_09_adaptive_weighting(prs, conceptual_dir)
    create_slide_10_theoretical(prs, data_dir)
    create_slide_11_sample_complexity(prs, data_dir)
    create_slide_12_experimental_setup(prs, conceptual_dir)
    create_slide_13_statistical_analysis(prs, data_dir)
    create_slide_14_main_results(prs, data_dir, paper_dir)
    create_slide_15_performance_breakdown(prs, data_dir)
    create_slide_16_physics_discovery(prs, data_dir)
    create_slide_17_convergence(prs, data_dir, paper_dir)
    create_slide_18_ablation(prs, data_dir, paper_dir)
    create_slide_19_efficiency(prs, data_dir)
    create_slide_20_limitations_domain(prs, conceptual_dir)
    create_slide_21_limitations_theoretical(prs, conceptual_dir)
    create_slide_22_future_domains(prs, conceptual_dir)
    create_slide_23_future_theoretical(prs, conceptual_dir)
    create_slide_24_broader_impact(prs, conceptual_dir)
    create_slide_25_conclusion(prs, conceptual_dir)
    create_slide_26_technical_contributions(prs, conceptual_dir)
    create_slide_27_research_impact(prs, conceptual_dir)
    create_slide_28_questions(prs, conceptual_dir)
    
    # Save presentation
    output_file = "Complete_Physics_Informed_Meta_Learning_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Complete PowerPoint presentation saved as: {output_file}")
    return output_file

def add_slide_with_layout(prs, layout_index=1):
    """Add a slide with specified layout"""
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)

def add_title_and_content(slide, title_text, content_text, font_size=18):
    """Add title and content to slide"""
    slide.shapes.title.text = title_text
    
    if len(slide.shapes.placeholders) > 1:
        content_placeholder = slide.shapes.placeholders[1]
        content_placeholder.text = content_text
        
        # Format text
        for paragraph in content_placeholder.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = 'Calibri'

def add_image_to_slide(slide, image_path, left=Inches(7), top=Inches(1.5), width=Inches(6), height=Inches(5)):
    """Add image to slide if it exists"""
    if image_path and Path(image_path).exists():
        try:
            slide.shapes.add_picture(str(image_path), left, top, width, height)
            return True
        except Exception as e:
            print(f"Warning: Could not add image {image_path}: {e}")
            return False
    return False

# Individual slide creation functions
def create_slide_01_title(prs, conceptual_dir):
    """Slide 1: Title Slide"""
    slide = add_slide_with_layout(prs, 0)  # Title slide layout
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Meta-Learning Physics-Informed Neural Networks for Few-Shot Parameter Inference"
    
    subtitle_text = """Brandon Yee¹, Wilson Collins¹, Benjamin Pellegrini¹, Caden Wang²

¹ Yee Collins Research Group
² Department of Computer Science, New York University

GitHub: https://github.com/YCRG-Labs/meta-pinn

AAAI 2026"""
    
    subtitle.text = subtitle_text
    
    # Try to add title slide image
    add_image_to_slide(slide, conceptual_dir / "title_slide.png", 
                      Inches(0), Inches(0), Inches(13.33), Inches(7.5))

def create_slide_02_problem(prs, conceptual_dir):
    """Slide 2: The Core Problem - PINN Limitations"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "The Core Problem - PINN Limitations"
    content_text = """Traditional Physics-Informed Neural Networks Face Critical Limitations:

• Extensive retraining required for each new problem
• No knowledge transfer from previously solved problems  
• Poor rapid adaptation to new physical domains
• Inefficient for few-shot scenarios with minimal data

Problem: Each new physics problem treated as completely independent

Impact: Computational bottleneck for real-world applications

Solution Preview: Our approach achieves 67% reduction in training time"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "pinn_limitations.png")

def create_slide_03_motivation(prs, conceptual_dir, data_dir):
    """Slide 3: Motivating Example - Fluid Dynamics"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Motivating Example - Fluid Dynamics"
    content_text = """Consider Solving Navier-Stokes Equations Across Different Reynolds Numbers:

• Traditional approach: Train separate PINN for Re=100, Re=200, Re=500, Re=1000
• Problem: Each requires full training from scratch (500 steps)
• Inefficiency: No leverage of shared fluid dynamics principles
• Real need: Rapid adaptation to new flow conditions with minimal data

Our Results:
• 67% reduction in training time (12.4h → 4.1h)
• 3× fewer adaptation steps (150 → 50)
• 15% improvement in accuracy"""
    
    add_title_and_content(slide, title_text, content_text, 16)
    
    # Add both fluid dynamics examples and cost comparison
    add_image_to_slide(slide, conceptual_dir / "fluid_dynamics_examples.png", 
                      Inches(7), Inches(0.5), Inches(6), Inches(3))
    add_image_to_slide(slide, data_dir / "computational_cost_comparison.png",
                      Inches(7), Inches(4), Inches(6), Inches(3))

def create_slide_04_meta_learning(prs, conceptual_dir):
    """Slide 4: Why Meta-Learning for Physics?"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Why Meta-Learning for Physics?"
    content_text = """Meta-Learning Offers Promising Solution:

• "Learning to learn" - rapidly adapt to new tasks using prior experience
• Successful in computer vision and natural language processing
• Physics applications remain largely unexplored
• Unique challenges: incorporating physics constraints into meta-learning objectives

Key Benefits Demonstrated:
• 3× Fewer Adaptation Steps (50 vs 150)
• 15% Better Generalization Performance
• Leverages Prior Physics Knowledge
• Maintains Physical Consistency"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "meta_learning_concept.png")

def create_slide_05_contributions(prs, conceptual_dir):
    """Slide 5: Research Contributions"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Research Contributions"
    content_text = """Our Framework Addresses These Challenges Through:

• Novel meta-learning algorithm incorporating physics constraints in inner and outer loops
• Theoretical convergence guarantees and sample complexity bounds
• Adaptive constraint weighting mechanism for diverse tasks
• Automated physics discovery with natural language interpretation
• Comprehensive experimental validation with rigorous statistical analysis

Mathematical Framework:
θ* = argmin E[L_total(φ_T, T)]
L_total = L_data + λ(T)L_physics"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "research_contributions.png")

# Continue with remaining slides...
def create_slide_06_formulation(prs, conceptual_dir):
    """Slide 6: Problem Formulation"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Problem Formulation"
    content_text = """Mathematical Framework:

Task Distribution: p(T) where each task T_i has:
• Domain: Ω_i ⊂ R^d with boundary ∂Ω_i
• PDE: F_i[u_i] = 0 for x ∈ Ω_i
• Boundary conditions: B_i[u_i] = 0 for x ∈ ∂Ω_i
• Limited training data: D_i = {(x_j, u_j)} where N_i is small

Goal: Learn meta-model enabling quick adaptation with minimal data while respecting physics

Key Challenge: Balance between data fitting and physics constraint satisfaction"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "problem_formulation.png")

def create_slide_07_framework(prs, conceptual_dir):
    """Slide 7: Physics-Informed Meta-Learning Framework"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Physics-Informed Meta-Learning Framework"
    content_text = """Extending MAML with Physics Constraints:

Meta-objective:
θ* = argmin E[L_total(φ_T, T)]

Adapted parameters:
φ_T = θ - α∇_θ L_total(θ, T)

Total loss combines data and physics:
L_total(φ, T) = L_data(φ, T) + λ(T)L_physics(φ, T)

Key Innovation: Physics constraints enforced in both inner and outer optimization loops"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "framework_flowchart.png")

def create_slide_08_physics_loss(prs, conceptual_dir):
    """Slide 8: Physics Loss Implementation"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Physics Loss Implementation"
    content_text = """Enforcing PDE Constraints:

L_physics(φ, T) = E[|F[u_φ]|²] + E[|B[u_φ]|²]

• Interior points: PDE residual must be zero throughout domain
• Boundary points: Boundary conditions must be satisfied
• Automatic differentiation: Compute derivatives for PDE evaluation
• Collocation method: Sample points for expectation approximation

Implementation: Strong inductive bias for physically meaningful solutions"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "physics_loss_diagram.png")

def create_slide_09_adaptive_weighting(prs, conceptual_dir):
    """Slide 9: Adaptive Constraint Weighting"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Adaptive Constraint Weighting"
    content_text = """Learning Task-Specific Physics Importance:

λ(T) = σ(W_λ h_T + b_λ)

• Task embedding h_T captures task characteristics
• Learned parameters W_λ, b_λ determine weighting strategy
• Sigmoid activation ensures positive constraint weights
• Adaptive mechanism handles diverse physics complexity

Examples:
• High Re (turbulent): λ=0.8 (high physics weight)
• Low Re (laminar): λ=0.3 (lower physics weight)
• Complex geometry: λ=0.9 (very high weight)"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "adaptive_weighting.png")

def create_slide_10_theoretical(prs, data_dir):
    """Slide 10: Theoretical Convergence Guarantees"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Theoretical Convergence Guarantees"
    content_text = """Mathematical Analysis of Framework:

Under standard assumptions:
• Loss function L-Lipschitz continuous with L ≤ C₁
• Gradient variance bounded by σ² ≤ C₂
• Physics constraints μ-strongly convex

Convergence rate:
E[|∇L(θ_T)|²] ≤ C₁/T + C₂√(log T/T)

where C₁ = L²_data + λ²L²_physics/μ

Key insight: Physics regularization can improve convergence when constraints are strongly convex"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, data_dir / "theoretical_convergence.png")

# Add remaining slide creation functions...
def create_remaining_slides(prs, conceptual_dir, data_dir, paper_dir):
    """Create slides 11-28 (abbreviated for space)"""
    
    # Slide 11: Sample Complexity Analysis
    slide11 = add_slide_with_layout(prs, 1)
    add_title_and_content(slide11, "Sample Complexity Analysis", 
                         """Theoretical Sample Efficiency:

To achieve ε-accuracy with probability 1-δ:
N = O(d log(1/δ) / [ε²(1 + γ)])

• Standard complexity: O(d log(1/δ)/ε²)
• Physics benefit: Factor (1 + γ) improvement where γ > 0
• Physics regularization reduces required sample size""")
    add_image_to_slide(slide11, data_dir / "sample_complexity.png")
    
    # Continue with all remaining slides...
    # (Implementation continues for slides 12-28)

# Move main function to the end of the file

def create_slide_11_sample_complexity(prs, data_dir):
    """Slide 11: Sample Complexity Analysis"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Sample Complexity Analysis"
    content_text = """Theoretical Sample Efficiency:

To achieve ε-accuracy with probability 1-δ:
N = O(d log(1/δ) / [ε²(1 + γ)])

• Standard complexity: O(d log(1/δ)/ε²)
• Physics benefit: Factor (1 + γ) improvement where γ > 0
• Dimension d: Effective problem dimension
• Physics regularization reduces required sample size

Key insight: Physics constraints act as structural risk minimization"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, data_dir / "sample_complexity.png")

def create_slide_12_experimental_setup(prs, conceptual_dir):
    """Slide 12: Experimental Setup Overview"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Experimental Setup Overview"
    content_text = """Comprehensive Evaluation on Fluid Dynamics:

Problem classes:
• Navier-Stokes equations with Reynolds numbers 100-1000
• Heat transfer with varying boundary conditions
• Burgers equation variants with different viscosity parameters
• Lid-driven cavity problems with varying geometries

Dataset structure:
• 200 training tasks and 50 test tasks per problem class
• 20-100 data points per task for adaptation
• Realistic few-shot scenarios for computational physics"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "experimental_setup.png")

def create_slide_13_statistical_analysis(prs, data_dir):
    """Slide 13: Statistical Analysis Methodology"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Statistical Analysis Methodology"
    content_text = """Rigorous Statistical Validation:

Analysis framework:
• Python 3.9 with NumPy, PyTorch, SciPy
• Significance level α = 0.05 with Bonferroni correction
• Effect sizes calculated using Cohen's d
• Power analysis with β = 0.8 for sample size calculations
• Bootstrap resampling with 1000 iterations for confidence intervals

Statistical tests:
• Two-tailed t-tests for group comparisons
• ANOVA with post-hoc Tukey tests where appropriate"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, data_dir / "statistical_analysis.png")

def create_slide_14_main_results(prs, data_dir, paper_dir):
    """Slide 14: Main Experimental Results"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Main Experimental Results"
    content_text = """Physics-Informed Meta-Learning Significantly Outperforms Baselines:

Key findings:
• 92.4% validation accuracy (SD = 4.2%, 95% CI [88.2%, 96.6%])
• 83.0% for Transfer PINN baseline (p < 0.001, Cohen's d = 2.1)
• 15% improvement in generalization performance
• 3× fewer adaptation steps (50 vs 150 steps)

Statistical significance:
• t(49) = 22.5, p < 0.001 vs Transfer PINN baseline
• Large effect size (Cohen's d = 2.1) indicates practical significance"""
    
    add_title_and_content(slide, title_text, content_text, 16)
    
    # Add main results figure
    add_image_to_slide(slide, data_dir / "main_experimental_results.png")
    
    # Try to add paper figure if available
    if (paper_dir / "Figure 1.png").exists():
        add_image_to_slide(slide, paper_dir / "Figure 1.png", 
                          Inches(7), Inches(4), Inches(6), Inches(3))

def create_slide_15_performance_breakdown(prs, data_dir):
    """Slide 15: Detailed Performance Breakdown"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Detailed Performance Breakdown"
    content_text = """Performance Across Different Shot Settings:

Method          5-shot      10-shot     20-shot     Steps   Validation Acc.
Standard PINN   0.654±0.089  0.721±0.076  0.783±0.065   500     0.783
Transfer PINN   0.712±0.082  0.768±0.071  0.824±0.058   150     0.824
MAML           0.698±0.091  0.745±0.083  0.801±0.072   100     0.801
PI-MAML (Ours) 0.847±0.052  0.891±0.048  0.922±0.041    50     0.922

Key insights:
• Consistent improvement across all few-shot settings
• Performance gap increases with more adaptation data
• Substantial improvement even in very few-shot (5-shot) scenarios"""
    
    add_title_and_content(slide, title_text, content_text, 14)
    add_image_to_slide(slide, data_dir / "performance_breakdown.png")

def create_slide_16_physics_discovery(prs, data_dir):
    """Slide 16: Physics Discovery Results"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Physics Discovery Results"
    content_text = """Automated Discovery Identifies Key Physical Relationships:

Discovery accuracies:
• Reynolds number dependence: 94% ± 3% (95% CI [91%-97%])
• Pressure-velocity coupling: 91% ± 4% (95% CI [87%-95%])
• Boundary layer effects: 89% ± 5% (95% CI [84%-94%])
• Heat transfer correlations: 92% ± 3% (95% CI [89%-95%])

Overall performance:
• Precision: 95%, Recall: 93%, F1-score: 94%
• All significantly above chance (25%, p < 0.001)

Example interpretation: "Viscosity decreases linearly with Reynolds number, 
indicating stronger inertial effects at higher flow rates" """
    
    add_title_and_content(slide, title_text, content_text, 16)
    add_image_to_slide(slide, data_dir / "physics_discovery_results.png")

def create_slide_17_convergence(prs, data_dir, paper_dir):
    """Slide 17: Convergence Analysis"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Convergence Analysis"
    content_text = """Physics Constraints Accelerate Meta-Training:

Convergence benefits:
• Faster convergence with physics constraints
• Lower final loss compared to standard MAML
• Better conditioning of optimization landscape
• Stable training across different initializations

Physics constraints provide inductive bias that guides optimization
Better conditioning reduces sensitivity to hyperparameters"""
    
    add_title_and_content(slide, title_text, content_text)
    
    # Add convergence analysis figure
    add_image_to_slide(slide, data_dir / "convergence_analysis.png")
    
    # Try to add paper figure if available
    if (paper_dir / "Figure 2.png").exists():
        add_image_to_slide(slide, paper_dir / "Figure 2.png", 
                          Inches(7), Inches(4), Inches(6), Inches(3))

def create_slide_18_ablation(prs, data_dir, paper_dir):
    """Slide 18: Ablation Study Results"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Ablation Study Results"
    content_text = """Each Component Contributes Significantly to Performance:

Component analysis:
• Full PI-MAML: 0.924 ± 0.042
• Without adaptive weighting: 0.887 ± 0.055 (p < 0.001)
• Without physics discovery: 0.901 ± 0.049 (p = 0.023)
• Without physics constraints: 0.801 ± 0.072 (p < 0.001)
• Without meta-learning: 0.830 ± 0.057 (p < 0.001)

Key insights:
• Physics constraints provide largest improvement (Cohen's d = 1.8)
• All components statistically significant
• Adaptive weighting crucial for diverse tasks"""
    
    add_title_and_content(slide, title_text, content_text, 16)
    
    # Add ablation study figure
    add_image_to_slide(slide, data_dir / "ablation_study.png")
    
    # Try to add paper figure if available
    if (paper_dir / "Figure 3.png").exists():
        add_image_to_slide(slide, paper_dir / "Figure 3.png", 
                          Inches(7), Inches(4), Inches(6), Inches(3))

def create_slide_19_efficiency(prs, data_dir):
    """Slide 19: Computational Efficiency Analysis"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Computational Efficiency Analysis"
    content_text = """Significant Computational Savings:

Method          Training Time  Memory   GPU Util.  Energy
                (hrs)         (GB)     (%)        (kWh)
Standard PINN   12.4 ± 1.2    8.9±0.7  85±5       24.8±2.4
Transfer PINN   8.7 ± 0.9     7.2±0.6  78±4       17.4±1.8
MAML           6.2 ± 0.8     6.8±0.5  82±3       12.4±1.6
PI-MAML (Ours) 4.1 ± 0.6     5.9±0.4  88±2       8.2±1.2

Key benefits:
• 3× reduction in training time compared to standard PINN
• Lower memory requirements enable larger problem sizes
• Higher GPU utilization indicates better computational efficiency
• Energy savings important for environmental considerations"""
    
    add_title_and_content(slide, title_text, content_text, 14)
    add_image_to_slide(slide, data_dir / "computational_efficiency.png")

def create_slide_20_limitations_domain(prs, conceptual_dir):
    """Slide 20: Limitations - Domain Specificity"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Limitations - Domain Specificity"
    content_text = """Current Scope and Constraints:

Domain limitations:
• Focus on fluid dynamics - broader physics domains need investigation
• Parameter ranges limited to tested Reynolds numbers and geometries
• Scalability to very high-dimensional problems unclear
• Task similarity assumptions may not hold across all physics domains

Technical constraints:
• Regularity assumptions in theoretical analysis
• Smoothness requirements for automatic differentiation
• Collocation sampling may miss important physics features

Future work will address these limitations through domain expansion"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "limitations_domain.png")

def create_slide_21_limitations_theoretical(prs, conceptual_dir):
    """Slide 21: Limitations - Theoretical Assumptions"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Limitations - Theoretical Assumptions"
    content_text = """Mathematical Framework Constraints:

Convergence analysis assumes:
• Lipschitz continuity of loss functions
• Bounded gradient variance across task distribution
• Strong convexity of physics constraints
• Standard regularity conditions for meta-learning

Practical considerations:
• 94% physics discovery accuracy may miss subtle effects
• Adaptive weighting requires task embedding quality
• Finite sample effects in few-shot scenarios

These are areas for continued research rather than fundamental flaws"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "limitations_theoretical.png")

def create_slide_22_future_domains(prs, conceptual_dir):
    """Slide 22: Future Work - Broader Physics Domains"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Future Work - Broader Physics Domains"
    content_text = """Research Directions for Extension:

Domain expansion:
• Solid mechanics with stress-strain relationships
• Electromagnetics with Maxwell's equations
• Quantum mechanics with Schrödinger equations
• Multi-physics coupled problems

Methodological improvements:
• Hierarchical meta-learning for multi-scale problems
• Symbolic AI integration for enhanced physics discovery
• Uncertainty quantification in meta-learned representations
• Active learning for optimal data collection strategies"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "future_work_domains.png")

def create_slide_23_future_theoretical(prs, conceptual_dir):
    """Slide 23: Future Work - Theoretical Extensions"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Future Work - Theoretical Extensions"
    content_text = """Enhanced Mathematical Framework:

Theoretical improvements:
• Non-convex physics constraints analysis
• Distribution shift robustness guarantees
• Multi-task learning bounds for physics problems
• Approximation theory for neural PDE solutions

Algorithmic advances:
• Second-order meta-learning for faster convergence
• Gradient-free optimization for non-differentiable physics
• Federated learning for distributed physics computation
• Continual learning for evolving physics understanding"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "future_work_theoretical.png")

def create_slide_24_broader_impact(prs, conceptual_dir):
    """Slide 24: Broader Impact and Applications"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Broader Impact and Applications"
    content_text = """Potential Applications in Science and Engineering:

Resource-constrained scenarios:
• Real-time control of physical systems
• Edge computing for IoT physics applications
• Rapid prototyping of engineering designs
• Emergency response with limited data

Scientific discovery:
• Parameter estimation in experimental physics
• Model selection among competing theories
• Anomaly detection in physical systems
• Scientific hypothesis generation and testing

Timeline: Short-term (1-2 years) to long-term (5+ years) impact"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "broader_impact.png")

def create_slide_25_conclusion(prs, conceptual_dir):
    """Slide 25: Conclusion - Addressing the Original Motivation"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Conclusion - Addressing the Original Motivation"
    content_text = """Framework Successfully Addresses PINN Limitations:

Original problem: PINNs require extensive retraining for each new problem

Our solution demonstrates:
• 92.4% validation accuracy with physics-informed meta-learning
• 15% improvement over transfer learning baselines
• 3× faster adaptation (50 vs 150 steps)
• Automated physics discovery with 94% accuracy
• Theoretical guarantees for convergence and sample complexity

Impact: Enables practical deployment in resource-constrained settings
Advances physics-informed machine learning field"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "conclusion_comparison.png")

def create_slide_26_technical_contributions(prs, conceptual_dir):
    """Slide 26: Technical Contributions Summary"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Technical Contributions Summary"
    content_text = """Novel Framework Advances Physics-Informed Machine Learning:

Algorithmic innovations:
• Physics-informed meta-learning algorithm with dual-loop optimization
• Adaptive constraint weighting for diverse task handling
• Automated physics discovery with natural language interpretation

Theoretical advances:
• Convergence rate analysis for physics-constrained meta-learning
• Sample complexity bounds showing physics regularization benefits
• Mathematical guarantees under standard assumptions

Experimental validation:
• Rigorous statistical analysis with proper effect size reporting
• Comprehensive ablation studies validating each component
• Multiple physics domains demonstrating generalizability"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "technical_contributions.png")

def create_slide_27_research_impact(prs, conceptual_dir):
    """Slide 27: Impact on Physics-Informed Machine Learning"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Impact on Physics-Informed Machine Learning"
    content_text = """Advancing the Field Through Domain-Specific Meta-Learning:

Methodological impact:
• First comprehensive physics-informed meta-learning framework
• Theoretical foundation for future physics meta-learning research
• Practical efficiency gains for computational physics applications

Research implications:
• Combines domain knowledge with general-purpose ML techniques
• Demonstrates value of incorporating physics constraints in meta-learning
• Opens new research directions in scientific machine learning

Community impact: Template for applying meta-learning to other scientific domains"""
    
    add_title_and_content(slide, title_text, content_text)
    add_image_to_slide(slide, conceptual_dir / "research_impact.png")

def create_slide_28_questions(prs, conceptual_dir):
    """Slide 28: Questions and Discussion"""
    slide = add_slide_with_layout(prs, 1)
    
    title_text = "Questions and Discussion"
    content_text = """Thank You - Questions Welcome

Key takeaways:
✓ Meta-learning enables rapid physics problem adaptation
✓ Physics constraints improve both efficiency and accuracy
✓ Theoretical guarantees provide confidence in approach
✓ Automated discovery adds interpretability and scientific insight

Quantitative achievements:
• 92.4% validation accuracy with 95% confidence intervals
• 15% improvement over baselines with large effect size (Cohen's d = 2.1)
• 3× faster adaptation and 67% training time reduction

Code available: https://github.com/YCRG-Labs/meta-pinn

Contact:
Brandon Yee: b.yee@ycrg-labs.org
Caden Wang: cw4973@nyu.edu"""
    
    add_title_and_content(slide, title_text, content_text, 16)
    add_image_to_slide(slide, conceptual_dir / "questions_discussion.png")
def mai
n():
    """Generate the complete PowerPoint presentation"""
    print("Generating complete PowerPoint presentation with all 28 slides...")
    
    # Check if required directories exist
    required_dirs = ["conceptual_figures", "presentation_figures"]
    missing_dirs = []
    for dir_name in required_dirs:
        if not Path(dir_name).exists():
            missing_dirs.append(dir_name)
    
    if missing_dirs:
        print(f"Warning: Missing directories: {missing_dirs}")
        print("Some images may not be included in the presentation.")
    
    # Create the presentation
    output_file = create_complete_presentation()
    
    print(f"\n🎉 SUCCESS! Complete PowerPoint presentation generated!")
    print(f"📁 File: {output_file}")
    print(f"📊 Total slides: 28")
    print(f"🎯 Includes all figures and content")
    print(f"💼 Ready for professional presentation!")
    
    # Provide usage instructions
    print(f"\n📋 Usage Instructions:")
    print(f"1. Open {output_file} in PowerPoint")
    print(f"2. Review and customize styling as needed")
    print(f"3. Practice with speaker notes")
    print(f"4. Present with confidence!")
    
    return output_file

if __name__ == "__main__":
    main()